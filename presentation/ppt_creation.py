from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from io import BytesIO

def create_ppt_from_content(content, main_title):
    prs = Presentation()
    
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    margin_left = Inches(1.5)
    margin_right = Inches(1.5)
    content_width = prs.slide_width - margin_left - margin_right
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.width = content_width
    title_placeholder.height = Inches(2)
    title_placeholder.left = margin_left
    title_placeholder.top = Inches(2.75)
    
    subtitle_placeholder.width = content_width
    subtitle_placeholder.height = Inches(1.0)
    subtitle_placeholder.left = margin_left
    subtitle_placeholder.top = Inches(3)
    
    # Capitalize first letters of all words in the main title
    title_text = main_title.title()
    
    title_placeholder.text = title_text
    subtitle_placeholder.text = "gen. by ai text-to-ppt"
    
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.name = 'Times New Roman'
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_placeholder.text_frame.paragraphs[0].font.name = 'Times New Roman'
    subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
    
    titles = []
    
    slides_content = content.strip().split('## Slide ')
    for slide_content in slides_content:
        if slide_content:
            parts = slide_content.split('\n', 1)
            if len(parts) < 2:
                continue
            title, body = parts[0].strip(), parts[1].strip()
            
            if ': ' not in title:
                print(f"Skipping invalid title format: {title}")
                continue
            
            title_text = title.split(': ', 1)[1]
            titles.append(title_text)
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]

            title_placeholder.width = content_width
            title_placeholder.height = Inches(1.0)
            title_placeholder.left = margin_left
            title_placeholder.top = Inches(0.5)
            
            content_placeholder.width = content_width
            content_placeholder.height = Inches(5.5)
            content_placeholder.left = margin_left
            content_placeholder.top = Inches(1.5)

            title_text_frame = title_placeholder.text_frame
            title_text_frame.clear()
            p = title_text_frame.add_paragraph()
            p.text = title_text
            p.font.size = Pt(32)
            p.font.name = 'Times New Roman'
            p.font.bold = True
            
            content_text_frame = content_placeholder.text_frame
            content_text_frame.clear()
            body_lines = body.split('\n')
            for line in body_lines:
                if line.strip():
                    p = content_text_frame.add_paragraph()
                    line = line.lstrip()
                    if '?' in line:
                        p.text = line.lstrip('*').strip()
                        p.font.bold = True
                    else:
                        p.text = line
                    p.font.size = Pt(20)
                    p.font.name = 'Times New Roman'
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.text = paragraph.text.replace('*','')

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("- ", "")
    
    outline_slide_layout = prs.slide_layouts[1]
    outline_slide = prs.slides.add_slide(outline_slide_layout)
    title_placeholder = outline_slide.shapes.title
    content_placeholder = outline_slide.placeholders[1]

    title_placeholder.width = content_width
    title_placeholder.height = Inches(1.0)
    title_placeholder.left = margin_left
    title_placeholder.top = Inches(0.5)

    content_placeholder.width = content_width
    content_placeholder.height = Inches(5.5)
    content_placeholder.left = margin_left
    content_placeholder.top = Inches(1.5)
    
    title_placeholder.text = "Outline"
    title_text_frame = title_placeholder.text_frame
    title_text_frame.paragraphs[0].font.size = Pt(32)
    title_text_frame.paragraphs[0].font.name = 'Times New Roman'
    title_text_frame.paragraphs[0].font.bold = True
    
    content_text_frame = content_placeholder.text_frame
    content_text_frame.clear()
    for title in titles:
        p = content_text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.name = 'Times New Roman'
    
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(1, slides[-1])

    thank_you_slide_layout = prs.slide_layouts[1]
    thank_you_slide = prs.slides.add_slide(thank_you_slide_layout)
    title_placeholder = thank_you_slide.shapes.title
    
    title_placeholder.width = content_width
    title_placeholder.height = Inches(2)
    title_placeholder.left = margin_left
    title_placeholder.top = Inches(2.75)
    
    title_placeholder.text = "Thank You"
    title_text_frame = title_placeholder.text_frame
    title_text_frame.paragraphs[0].font.size = Pt(44)
    title_text_frame.paragraphs[0].font.name = 'Times New Roman'
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)  # Reset the buffer position to the beginning
    return pptx_bytes
